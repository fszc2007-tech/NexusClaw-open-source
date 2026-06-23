from __future__ import annotations

import csv
import json
import mimetypes
import os
import re
import subprocess
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

import httpx

from app.core.config import settings


SUPPORTED_NATIVE_EXTENSIONS = {
    "txt",
    "md",
    "csv",
    "html",
    "htm",
    "docx",
    "xlsx",
    "xlsm",
    "pptx",
    "pptm",
}
SUPPORTED_OCR_EXTENSIONS = {
    "png",
    "jpg",
    "jpeg",
    "webp",
    "bmp",
    "tif",
    "tiff",
}
LEGACY_OFFICE_EXTENSIONS = {"doc", "xls", "ppt"}


@dataclass
class ParsedBlock:
    block_type: str
    text: str
    order_no: int
    page_no: int | None = None
    sheet_name: str | None = None
    slide_no: int | None = None
    section_title: str | None = None
    metadata: dict[str, Any] = field(default_factory=dict)

    def to_dict(self) -> dict[str, Any]:
        return {
            "block_type": self.block_type,
            "text": self.text,
            "order_no": self.order_no,
            "page_no": self.page_no,
            "sheet_name": self.sheet_name,
            "slide_no": self.slide_no,
            "section_title": self.section_title,
            "metadata": self.metadata,
        }


@dataclass
class ParsedDocument:
    parser_name: str
    text: str
    blocks: list[ParsedBlock]
    metadata: dict[str, Any] = field(default_factory=dict)

    def to_dict(self) -> dict[str, Any]:
        return {
            "parser_name": self.parser_name,
            "text": self.text,
            "metadata": self.metadata,
            "blocks": [block.to_dict() for block in self.blocks],
        }


class DocumentParser:
    def __init__(self) -> None:
        self.project_root = Path(__file__).resolve().parents[2]

    def parse(
        self,
        file_path: Path,
        file_ext: str | None,
        mime_type: str | None,
    ) -> tuple[ParsedDocument | None, str | None]:
        normalized_ext = (file_ext or file_path.suffix.lstrip(".")).lower()
        normalized_mime = mime_type or mimetypes.guess_type(file_path.name)[0]
        normalized_ext = normalized_ext or self._guess_extension_from_mime(normalized_mime)

        try:
            if normalized_ext in LEGACY_OFFICE_EXTENSIONS:
                return None, f"暂不支持 .{normalized_ext} 老格式，请转换为 docx/xlsx/pptx 后重试"
            if normalized_ext == "txt":
                return self._parse_text_document(file_path, normalized_ext), None
            if normalized_ext == "md":
                return self._parse_markdown_document(file_path), None
            if normalized_ext == "csv":
                return self._parse_csv_document(file_path), None
            if normalized_ext in {"html", "htm"}:
                return self._parse_html_document(file_path), None
            if normalized_ext == "docx":
                return self._parse_docx_document(file_path), None
            if normalized_ext in {"xlsx", "xlsm"}:
                return self._parse_xlsx_document(file_path), None
            if normalized_ext in {"pptx", "pptm"}:
                return self._parse_pptx_document(file_path), None
            if normalized_ext == "pdf":
                return self._parse_pdf_document(file_path), None
            if normalized_ext in SUPPORTED_OCR_EXTENSIONS or (normalized_mime or "").startswith("image/"):
                return self._parse_with_ocr(file_path, normalized_ext, normalized_mime, strategy="ocr_image"), None
            return None, "该格式暂未支持自动解析"
        except UnicodeDecodeError:
            return None, "文件编码无法识别，请转换为 UTF-8 或 GB18030 后重试"
        except Exception as exc:  # noqa: BLE001
            return None, f"解析失败: {exc}"

    def _guess_extension_from_mime(self, mime_type: str | None) -> str:
        mime_to_ext = {
            "text/plain": "txt",
            "text/markdown": "md",
            "text/csv": "csv",
            "text/html": "html",
            "application/pdf": "pdf",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
        }
        if not mime_type:
            return ""
        return mime_to_ext.get(mime_type, "")

    def _parse_text_document(self, file_path: Path, file_ext: str) -> ParsedDocument:
        text = _read_text_with_fallback(file_path)
        blocks = [
            ParsedBlock(
                block_type="paragraph",
                text=paragraph,
                order_no=index,
            )
            for index, paragraph in enumerate(_split_paragraphs(text), start=1)
        ]
        return ParsedDocument(
            parser_name=f"native_{file_ext}",
            text=text,
            blocks=blocks,
            metadata={
                "block_count": len(blocks),
                "route_kind": "native_text",
                "chunk_strategy": "structured_text",
                "file_category": "text",
            },
        )

    def _parse_markdown_document(self, file_path: Path) -> ParsedDocument:
        text = _read_text_with_fallback(file_path)
        blocks: list[ParsedBlock] = []
        order_no = 1
        paragraph_lines: list[str] = []

        def flush_paragraph() -> None:
            nonlocal order_no
            if not paragraph_lines:
                return
            paragraph_text = "\n".join(paragraph_lines).strip()
            if paragraph_text:
                blocks.append(
                    ParsedBlock(
                        block_type="paragraph",
                        text=paragraph_text,
                        order_no=order_no,
                    )
                )
                order_no += 1
            paragraph_lines.clear()

        for raw_line in text.splitlines():
            line = raw_line.rstrip()
            if not line.strip():
                flush_paragraph()
                continue
            heading_match = re.match(r"^(#{1,6})\s+(.*)$", line)
            if heading_match:
                flush_paragraph()
                level = len(heading_match.group(1))
                title = heading_match.group(2).strip()
                if title:
                    blocks.append(
                        ParsedBlock(
                            block_type="title",
                            text=title,
                            order_no=order_no,
                            section_title=title,
                            metadata={"heading_level": level},
                        )
                    )
                    order_no += 1
                continue
            paragraph_lines.append(line)

        flush_paragraph()
        return ParsedDocument(
            parser_name="native_md",
            text=text,
            blocks=blocks,
            metadata={
                "block_count": len(blocks),
                "route_kind": "native_text",
                "chunk_strategy": "structured_text",
                "file_category": "markdown",
            },
        )

    def _parse_csv_document(self, file_path: Path) -> ParsedDocument:
        rows: list[list[str]] = []
        csv_text = _read_text_with_fallback(file_path)
        for row in csv.reader(csv_text.splitlines()):
            normalized = [str(cell).strip() for cell in row]
            if any(normalized):
                rows.append(normalized)
        table_lines = [" | ".join(row) for row in rows]
        text = "\n".join(table_lines)
        header_fields = rows[0] if rows else []
        blocks = [
            ParsedBlock(
                block_type="table_row",
                text=line,
                order_no=index,
                metadata={"row_index": index, "header_fields": header_fields if index > 1 else normalized_header(row=line)},
            )
            for index, line in enumerate(table_lines, start=1)
        ]
        return ParsedDocument(
            parser_name="native_csv",
            text=text,
            blocks=blocks,
            metadata={
                "row_count": len(rows),
                "block_count": len(blocks),
                "route_kind": "native_table",
                "chunk_strategy": "spreadsheet",
                "file_category": "csv",
                "header_fields": header_fields,
            },
        )

    def _parse_html_document(self, file_path: Path) -> ParsedDocument:
        try:
            from bs4 import BeautifulSoup
        except ImportError as exc:  # pragma: no cover
            raise RuntimeError("缺少 beautifulsoup4 依赖，请先安装后重试") from exc

        html = _read_text_with_fallback(file_path)
        soup = BeautifulSoup(html, "html.parser")
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()

        blocks: list[ParsedBlock] = []
        order_no = 1
        for node in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "table"]):
            text = " ".join(node.get_text(" ", strip=True).split())
            if not text:
                continue
            tag_name = node.name or "p"
            block_type = "title" if tag_name.startswith("h") else "table" if tag_name == "table" else "paragraph"
            blocks.append(
                ParsedBlock(
                    block_type=block_type,
                    text=text,
                    order_no=order_no,
                    section_title=text if block_type == "title" else None,
                    metadata={"heading_level": int(tag_name[1]) if block_type == "title" and len(tag_name) == 2 else None},
                )
            )
            order_no += 1

        if not blocks:
            fallback_text = " ".join(soup.get_text(" ", strip=True).split())
            blocks.append(ParsedBlock(block_type="paragraph", text=fallback_text, order_no=1))

        text = "\n".join(block.text for block in blocks)
        return ParsedDocument(
            parser_name="native_html",
            text=text,
            blocks=blocks,
            metadata={
                "block_count": len(blocks),
                "title": (soup.title.string.strip() if soup.title and soup.title.string else None),
                "route_kind": "native_html",
                "chunk_strategy": "structured_text",
                "file_category": "html",
            },
        )

    def _parse_docx_document(self, file_path: Path) -> ParsedDocument:
        try:
            from docx import Document
        except ImportError as exc:  # pragma: no cover
            raise RuntimeError("缺少 python-docx 依赖，请先安装后重试") from exc

        document = Document(str(file_path))
        blocks: list[ParsedBlock] = []
        order_no = 1

        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            style_name = paragraph.style.name.lower() if paragraph.style and paragraph.style.name else ""
            block_type = "title" if "heading" in style_name or "标题" in style_name else "paragraph"
            heading_level = self._extract_docx_heading_level(style_name) if block_type == "title" else None
            blocks.append(
                ParsedBlock(
                    block_type=block_type,
                    text=text,
                    order_no=order_no,
                    section_title=text if block_type == "title" else None,
                    metadata={"heading_level": heading_level},
                )
            )
            order_no += 1

        for table_index, table in enumerate(document.tables, start=1):
            row_texts: list[str] = []
            for row in table.rows:
                cells = [" ".join(cell.text.split()) for cell in row.cells]
                if any(cells):
                    row_texts.append(" | ".join(cells))
            if not row_texts:
                continue
            blocks.append(
                ParsedBlock(
                    block_type="table",
                    text="\n".join(row_texts),
                    order_no=order_no,
                    metadata={"table_index": table_index, "row_count": len(row_texts)},
                )
            )
            order_no += 1

        text = "\n".join(block.text for block in blocks)
        return ParsedDocument(
            parser_name="native_docx",
            text=text,
            blocks=blocks,
            metadata={
                "block_count": len(blocks),
                "table_count": len(document.tables),
                "route_kind": "native_office",
                "chunk_strategy": "structured_text",
                "file_category": "docx",
            },
        )

    def _parse_xlsx_document(self, file_path: Path) -> ParsedDocument:
        try:
            from openpyxl import load_workbook
        except ImportError as exc:  # pragma: no cover
            raise RuntimeError("缺少 openpyxl 依赖，请先安装后重试") from exc

        workbook = load_workbook(filename=str(file_path), read_only=True, data_only=True)
        blocks: list[ParsedBlock] = []
        order_no = 1
        for sheet in workbook.worksheets:
            blocks.append(
                ParsedBlock(
                    block_type="title",
                    text=f"Sheet: {sheet.title}",
                    order_no=order_no,
                    sheet_name=sheet.title,
                    section_title=sheet.title,
                    metadata={"heading_level": 1},
                )
            )
            order_no += 1

            header_fields: list[str] | None = None
            for row_index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                cells = [str(value).strip() for value in row if value not in (None, "")]
                if not cells:
                    continue
                if header_fields is None:
                    header_fields = cells
                line = " | ".join(cells)
                blocks.append(
                    ParsedBlock(
                        block_type="table_row",
                        text=line,
                        order_no=order_no,
                        sheet_name=sheet.title,
                        metadata={"row_index": row_index, "header_fields": header_fields},
                    )
                )
                order_no += 1

        text = "\n".join(block.text for block in blocks)
        return ParsedDocument(
            parser_name="native_xlsx",
            text=text,
            blocks=blocks,
            metadata={
                "sheet_count": len(workbook.worksheets),
                "block_count": len(blocks),
                "route_kind": "native_office",
                "chunk_strategy": "spreadsheet",
                "file_category": "xlsx",
            },
        )

    def _parse_pptx_document(self, file_path: Path) -> ParsedDocument:
        try:
            from pptx import Presentation
        except ImportError as exc:  # pragma: no cover
            raise RuntimeError("缺少 python-pptx 依赖，请先安装后重试") from exc

        presentation = Presentation(str(file_path))
        blocks: list[ParsedBlock] = []
        order_no = 1
        for slide_no, slide in enumerate(presentation.slides, start=1):
            title_shape = slide.shapes.title
            if title_shape and getattr(title_shape, "text", "").strip():
                title_text = title_shape.text.strip()
                blocks.append(
                    ParsedBlock(
                        block_type="title",
                        text=title_text,
                        order_no=order_no,
                        slide_no=slide_no,
                        section_title=title_text,
                    )
                )
                order_no += 1

            for shape in slide.shapes:
                if title_shape is not None and shape == title_shape:
                    continue
                text = getattr(shape, "text", None)
                if text:
                    normalized = " ".join(text.split())
                    if normalized:
                        blocks.append(
                            ParsedBlock(
                                block_type="paragraph",
                                text=normalized,
                                order_no=order_no,
                                slide_no=slide_no,
                                section_title=title_shape.text.strip() if title_shape and getattr(title_shape, "text", "").strip() else None,
                            )
                        )
                        order_no += 1
                    continue

                if not getattr(shape, "has_table", False):
                    continue
                rows: list[str] = []
                for row in shape.table.rows:
                    cells = [" ".join(cell.text.split()) for cell in row.cells]
                    if any(cells):
                        rows.append(" | ".join(cells))
                if rows:
                    blocks.append(
                        ParsedBlock(
                            block_type="table",
                            text="\n".join(rows),
                            order_no=order_no,
                            slide_no=slide_no,
                        )
                    )
                    order_no += 1

        text = "\n".join(block.text for block in blocks)
        return ParsedDocument(
            parser_name="native_pptx",
            text=text,
            blocks=blocks,
            metadata={
                "slide_count": len(presentation.slides),
                "block_count": len(blocks),
                "route_kind": "native_office",
                "chunk_strategy": "slide",
                "file_category": "pptx",
            },
        )

    def _parse_pdf_document(self, file_path: Path) -> ParsedDocument:
        pdf_parse_mode = (settings.FILE_PDF_PARSE_MODE or "ocr_first").lower()
        if pdf_parse_mode == "ocr_only":
            return self._parse_with_ocr(file_path, "pdf", "application/pdf", strategy="ocr_pdf")
        if pdf_parse_mode == "ocr_first":
            try:
                return self._parse_with_ocr(file_path, "pdf", "application/pdf", strategy="ocr_pdf")
            except Exception:
                pass

        try:
            from pypdf import PdfReader
        except ImportError as exc:  # pragma: no cover
            raise RuntimeError("缺少 pypdf 依赖，请先安装后重试") from exc

        reader = PdfReader(str(file_path))
        if reader.is_encrypted:
            raise RuntimeError("PDF 已加密，暂不支持自动解析")

        blocks: list[ParsedBlock] = []
        order_no = 1
        texts: list[str] = []
        for page_no, page in enumerate(reader.pages, start=1):
            page_text = (page.extract_text() or "").strip()
            if not page_text:
                continue
            texts.append(page_text)
            blocks.append(
                ParsedBlock(
                    block_type="paragraph",
                    text=page_text,
                    order_no=order_no,
                    page_no=page_no,
                    metadata={"page_index": page_no},
                )
            )
            order_no += 1

        text = "\n\n".join(texts)
        if len(text.strip()) < settings.FILE_PDF_MIN_TEXT_CHARS_FOR_NATIVE:
            return self._parse_with_ocr(file_path, "pdf", "application/pdf", strategy="ocr_pdf")

        return ParsedDocument(
            parser_name="native_pdf",
            text=text,
            blocks=blocks,
            metadata={
                "page_count": len(reader.pages),
                "block_count": len(blocks),
                "route_kind": "native_pdf",
                "chunk_strategy": "paged_document",
                "file_category": "pdf",
            },
        )

    def _parse_with_ocr(
        self,
        file_path: Path,
        file_ext: str | None,
        mime_type: str | None,
        strategy: str,
    ) -> ParsedDocument:
        backend = (settings.PADDLE_OCR_BACKEND or "auto").lower()
        if backend in {"auto", "local_paddle"}:
            try:
                return self._parse_with_local_paddle_ocr(file_path)
            except Exception:
                if backend == "local_paddle":
                    raise
        if not settings.PADDLE_OCR_PARSE_URL:
            target_name = "本地 PaddleOCR 或远程 PaddleOCR-VL 服务"
            raise RuntimeError(f"当前文件需要 OCR 解析，但未配置 {target_name}")

        with file_path.open("rb") as file_handle:
            files = {
                "file": (
                    file_path.name,
                    file_handle,
                    mime_type or mimetypes.guess_type(file_path.name)[0] or "application/octet-stream",
                )
            }
            data = {
                "strategy": strategy,
                "file_ext": file_ext or "",
                "response_format": "json",
            }
            headers = {}
            if settings.PADDLE_OCR_API_KEY:
                headers["Authorization"] = f"Bearer {settings.PADDLE_OCR_API_KEY}"

            response = httpx.post(
                settings.PADDLE_OCR_PARSE_URL,
                files=files,
                data=data,
                headers=headers,
                timeout=settings.PADDLE_OCR_TIMEOUT_SECONDS,
            )

        response.raise_for_status()
        payload = response.json()
        blocks_payload = payload.get("blocks") or payload.get("elements") or []
        normalized_blocks: list[ParsedBlock] = []
        for index, item in enumerate(blocks_payload, start=1):
            text = str(item.get("text") or item.get("markdown") or item.get("content") or "").strip()
            if not text:
                continue
            normalized_blocks.append(
                ParsedBlock(
                    block_type=str(item.get("block_type") or item.get("type") or "paragraph"),
                    text=text,
                    order_no=index,
                    page_no=item.get("page_no"),
                    sheet_name=item.get("sheet_name"),
                    slide_no=item.get("slide_no"),
                    section_title=item.get("section_title"),
                    metadata={k: v for k, v in item.items() if k not in {"text", "markdown", "content", "type", "block_type", "page_no", "sheet_name", "slide_no", "section_title"}},
                )
            )

        text = str(payload.get("markdown") or payload.get("text") or payload.get("content") or "").strip()
        if not text and normalized_blocks:
            text = "\n".join(block.text for block in normalized_blocks)
        if not text:
            raise RuntimeError("OCR 服务未返回可用文本")

        if not normalized_blocks:
            normalized_blocks = [
                ParsedBlock(block_type="paragraph", text=paragraph, order_no=index)
                for index, paragraph in enumerate(_split_paragraphs(text), start=1)
            ]

        return ParsedDocument(
            parser_name=payload.get("parser_name") or "paddleocr_vl_1_5",
            text=text,
            blocks=normalized_blocks,
            metadata={
                "ocr_strategy": strategy,
                "page_count": payload.get("page_count"),
                "block_count": len(normalized_blocks),
                "provider": payload.get("provider") or "PaddleOCR-VL-1.5",
                "route_kind": "ocr_remote",
                "chunk_strategy": "ocr_document",
                "file_category": "visual_document",
            },
        )

    def _parse_with_local_paddle_ocr(self, file_path: Path) -> ParsedDocument:
        python_bin = self._resolve_runtime_path(settings.PADDLE_OCR_LOCAL_PYTHON_BIN)
        script_path = self._resolve_runtime_path(settings.PADDLE_OCR_LOCAL_SCRIPT_PATH)
        if not python_bin.exists():
            raise RuntimeError(f"未找到本地 OCR Python 解释器: {python_bin}")
        if not script_path.exists():
            raise RuntimeError(f"未找到本地 OCR 脚本: {script_path}")

        command = [
            str(python_bin),
            str(script_path),
            "--file",
            str(file_path.resolve()),
            "--lang",
            settings.PADDLE_OCR_LOCAL_LANG,
        ]
        det_model_dir = self._expand_optional_path(settings.PADDLE_OCR_LOCAL_DET_MODEL_DIR)
        rec_model_dir = self._expand_optional_path(settings.PADDLE_OCR_LOCAL_REC_MODEL_DIR)
        if det_model_dir:
            command.extend(["--det-model-dir", str(det_model_dir)])
        if rec_model_dir:
            command.extend(["--rec-model-dir", str(rec_model_dir)])

        env = {
            **os.environ,
            "PADDLE_PDX_DISABLE_MODEL_SOURCE_CHECK": "True",
        }
        result = subprocess.run(
            command,
            capture_output=True,
            text=True,
            timeout=settings.PADDLE_OCR_TIMEOUT_SECONDS,
            env=env,
            cwd=str(self.project_root),
            check=False,
        )
        if result.returncode != 0:
            error_message = (result.stderr or result.stdout).strip() or f"本地 OCR 进程退出码 {result.returncode}"
            raise RuntimeError(error_message)

        payload = json.loads(result.stdout)
        if payload.get("error"):
            raise RuntimeError(str(payload["error"]))
        blocks_payload = payload.get("blocks") or []
        normalized_blocks = [
            ParsedBlock(
                block_type=str(item.get("block_type") or "paragraph"),
                text=str(item.get("text") or "").strip(),
                order_no=int(item.get("order_no") or index),
                page_no=item.get("page_no"),
                sheet_name=item.get("sheet_name"),
                slide_no=item.get("slide_no"),
                section_title=item.get("section_title"),
                metadata=item.get("metadata") or {},
            )
            for index, item in enumerate(blocks_payload, start=1)
            if str(item.get("text") or "").strip()
        ]
        text = str(payload.get("text") or "").strip()
        if not text and normalized_blocks:
            text = "\n".join(block.text for block in normalized_blocks)
        if not text:
            raise RuntimeError("本地 OCR 未返回可用文本")
        return ParsedDocument(
            parser_name=str(payload.get("parser_name") or "local_paddle_ocr"),
            text=text,
            blocks=normalized_blocks,
            metadata={
                "ocr_strategy": "local_paddle",
                "page_count": payload.get("page_count"),
                "block_count": len(normalized_blocks),
                "provider": payload.get("provider") or "PaddleOCR",
                "route_kind": "ocr_local",
                "chunk_strategy": "ocr_document",
                "file_category": "visual_document",
            },
        )

    def _extract_docx_heading_level(self, style_name: str) -> int:
        match = re.search(r"(\d+)", style_name)
        if match:
            return max(1, min(6, int(match.group(1))))
        return 1

    def _resolve_runtime_path(self, raw_path: str | None) -> Path:
        if not raw_path:
            return self.project_root
        path = Path(raw_path).expanduser()
        if path.is_absolute():
            return path
        return self.project_root / path

    def _expand_optional_path(self, raw_path: str | None) -> Path | None:
        if not raw_path:
            return None
        path = Path(raw_path).expanduser()
        return path if path.exists() else None


def _read_text_with_fallback(file_path: Path) -> str:
    raw_bytes = file_path.read_bytes()
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "gbk"):
        try:
            return raw_bytes.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw_bytes.decode("utf-8", errors="ignore")


def _split_paragraphs(text: str) -> list[str]:
    paragraphs = [" ".join(line.split()) for line in text.splitlines() if line.strip()]
    return paragraphs or [text.strip()]


def normalized_header(row: str) -> list[str]:
    return [cell.strip() for cell in row.split("|") if cell.strip()]
